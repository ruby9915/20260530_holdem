# -*- coding: utf-8 -*-
"""
발표용 PPTX 생성 스크립트 (10~12장, 국문)
주제: 고정 상대 하 Tabular On-policy Monte Carlo 에서의 커버리지 천장 실증
실행: python make_slides.py
출력: 논문화/발표_커버리지천장.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT = "Malgun Gothic"          # 윈도우 기본 한글 폰트
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xEE, 0xF3, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_font(run, size=18, bold=False, color=GRAY, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def title_bar(slide, text, sub=None):
    # 상단 네이비 바
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    set_font(r, 26, True, RGBColor(0xFF, 0xFF, 0xFF))
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        set_font(r2, 14, False, RGBColor(0xCF, 0xDD, 0xEC))


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
            width=Inches(12.0), height=Inches(5.5), size=18, gap=10):
    tf = add_textbox(slide, left, top, width, height)
    first = True
    for it in items:
        lvl = int(it[0]) if isinstance(it, tuple) else 0
        txt = it[1] if isinstance(it, tuple) else it
        color = it[2] if (isinstance(it, tuple) and len(it) > 2 and it[2] is not None) else GRAY
        bold = it[3] if (isinstance(it, tuple) and len(it) > 3) else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = lvl
        mark = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark + txt
        set_font(r, size - lvl * 2, bold, color)


# ---------------------------------------------------------------------------
# Slide 1 — 표지
# ---------------------------------------------------------------------------
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tf = add_textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "고정 상대 하 Tabular On-policy 몬테카를로에서의\n커버리지 천장 실증"
set_font(r, 34, True, RGBColor(0xFF, 0xFF, 0xFF))
p2 = tf.add_paragraph(); p2.space_before = Pt(18)
r2 = p2.add_run()
r2.text = "탐색을 강화해도 천장은 움직이지 않는다 — 환경(고정 상대)이 천장을 결정한다"
set_font(r2, 18, False, RGBColor(0xCF, 0xDD, 0xEC))
p3 = tf.add_paragraph(); p3.space_before = Pt(24)
r3 = p3.add_run(); r3.text = "강화학습 기말 발표  ·  2026"
set_font(r3, 15, False, RGBColor(0x9F, 0xB6, 0xCE))
s.notes_slide.notes_text_frame.text = (
    "핵심 한 줄: 시험한 레버(예산·탐색·커버리지) 중에서는 상대가 만드는 전이 분포만이 "
    "성능 천장을 움직였다. tabular·on-policy MC 한정 조건부 명제임을 처음부터 명확히 한다."
)

# ---------------------------------------------------------------------------
# Slide 2 — 한 줄 요약 / 핵심 주장
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "핵심 주장 한 줄", "이 발표의 모든 슬라이드는 이 문장으로 수렴한다")
box = s.shapes.add_shape(1, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.9))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT; box.line.color.rgb = ACCENT
box.line.width = Pt(1.5)
tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = ("“탐색을 키워도 커버리지 천장은 불변하고, 고정 상대(환경)가 천장을 결정한다.”")
set_font(r, 22, True, NAVY)
bullets(s, [
    ("0", "단, tabular · on-policy Monte Carlo 환경에 한정된 조건부 명제", ACCENT, True),
    ("0", "‘유일하게 결정한다’는 강한 명제는 본 연구 범위 밖 → 정직한 한계로 명시"),
], top=Inches(4.4), size=18)
s.notes_slide.notes_text_frame.text = (
    "여기서 과장하지 않는다는 점을 먼저 못박는다. 조건부 명제, 경계 조건 명료화 + "
    "실증적·해석적 신규성이 우리의 기여다."
)

# ---------------------------------------------------------------------------
# Slide 3 — 문제 제기
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "1. 문제 제기 — 탐색 실패는 정말 ‘에이전트 탓’인가?")
bullets(s, [
    ("0", "통념: 탐색이 실패하면 원인을 학습자 측 결함으로 귀속", None, True),
    ("1", "탐색 알고리즘의 한계 / 데이터 커버리지 부족 / 신경망 함수 근사 오차"),
    ("0", "그러나 불완전정보 게임에서는 다른 일이 벌어질 수 있다", None, True),
    ("1", "상대의 은닉된 의사결정이 상태 전이 P(s'|s,a)를 직접 좌우"),
    ("1", "환경(상대)이 특정 상태로의 전이 경로를 원천적으로 닫으면,"),
    ("1", "학습자가 탐색을 아무리 키워도 그 상태 도달이 물리적으로 불가능"),
    ("0", "질문: 도달 실패가 ‘알고리즘 탓’인가, ‘환경의 구조적 차단’인가?", ACCENT, True),
], top=Inches(1.5), gap=12)
s.notes_slide.notes_text_frame.text = (
    "이 질문이 연구 전체의 출발점. 다음 슬라이드에서 둘을 분리하기 위한 실험 설계를 보여준다."
)

# ===========================================================================
# [배경 개념] 학부생 대상 기초 직관 슬라이드 (실험 환경 세팅 질문 대비)
# ===========================================================================

# --- 배경 0: 로드맵 / 용어 미리보기 ---
s = add_slide()
title_bar(s, "잠깐 — 오늘 나올 용어 미리 보기", "‘그게 뭔가요?’가 나올 단어를 먼저 한 줄씩 정리")
bullets(s, [
    ("0", "HUNL 포커: 1:1로 겨루는 노리밋 텍사스 홀덤 (우리의 실험장)", None, True),
    ("0", "상태/컨텍스트: 게임의 ‘상황’을 분류한 상자", None, True),
    ("0", "추상화: 거의 무한한 상황을 256개 상자로 압축", None, True),
    ("0", "고정 상대(페르소나): 학습하지 않는 정해진 스타일의 상대 = 환경", None, True),
    ("0", "도달 분포 d^π(s): 각 상자에 ‘얼마나 자주 가는지’", None, True),
    ("0", "몬테카를로 학습: 한 판을 끝까지 가보고 그 결과로 배움", None, True),
    ("0", "mbb/g: 한 판당 얼마나 따고 잃는지 재는 단위", None, True),
    ("1", "→ 지금부터 이 7개를 하나씩 그림으로 설명합니다", ACCENT, True),
], top=Inches(1.45), gap=10, size=17)
s.notes_slide.notes_text_frame.text = (
    "청중이 초반에 질문으로 끊지 않도록, 핵심 용어를 먼저 한 줄씩 보여주고 "
    "‘이제 하나씩 설명하겠다’고 안내한다."
)

# --- 배경 1: 포커 기초 (HUNL) ---
s = add_slide()
title_bar(s, "0-1. 실험장 — HUNL 포커가 뭔가요?")
bullets(s, [
    ("0", "HUNL = Heads-Up No-Limit Hold’em (헤즈업 노리밋 홀덤)", NAVY, True),
    ("1", "헤즈업: 딱 두 명이 1:1로 대결 (멀티플레이어 아님 → 분석이 깨끗)"),
    ("1", "노리밋: 베팅 금액에 상한이 없음 (올인까지 가능)"),
    ("0", "한 판의 흐름: 4개의 라운드", None, True),
    ("1", "Preflop → Flop → Turn → River (카드가 점점 공개됨)"),
    ("1", "각 라운드에서 폴드/체크/콜/레이즈 중 선택"),
    ("0", "왜 포커?  상대의 패가 ‘숨겨져’ 있는 불완전정보 게임", ACCENT, True),
    ("1", "시작 스택 200, 작은블라인드(SB)=1 / 큰블라인드(BB)=2, pokerkit 엔진 사용"),
], top=Inches(1.45), gap=10, size=17)
s.notes_slide.notes_text_frame.text = (
    "체스/바둑과 달리 상대 패가 안 보인다는 점이 핵심. 그래서 상대 행동이 정보를 가린 채 "
    "환경을 바꾼다. 스택/블라인드는 그냥 게임 세팅값."
)

# --- 배경 2: 왜 불완전정보가 중요한가 ---
s = add_slide()
title_bar(s, "0-2. 왜 ‘불완전정보’가 이 연구의 열쇠인가")
bullets(s, [
    ("0", "체스·바둑: 판이 다 보임 → 내가 탐색하면 어디든 갈 수 있음", None, True),
    ("0", "포커: 상대 패와 의도가 숨겨져 있음", RED, True),
    ("1", "상대가 특정 행동(예: 항상 폴드)을 하면 어떤 상황은 ‘아예 안 생김’"),
    ("1", "내가 아무리 다양하게 시도해도 그 상황엔 도달 불가"),
    ("0", "즉, 상대(=환경)가 ‘갈 수 있는 상황의 범위’를 정해버림", ACCENT, True),
    ("1", "이 발표의 핵심 현상: 탐색이 아니라 상대가 천장을 만든다"),
], top=Inches(1.5), gap=12, size=18)
s.notes_slide.notes_text_frame.text = (
    "여기서 ‘상대가 환경’이라는 개념을 직관적으로 심는다. 뒤의 모든 실험이 이 한 문장의 증명."
)

# --- 배경 3: 상태/컨텍스트/추상화 직관 ---
s = add_slide()
title_bar(s, "0-3. 상태·컨텍스트·추상화 — 상황을 ‘상자’에 넣기")
bullets(s, [
    ("0", "문제: 포커의 실제 상황 수는 천문학적 (카드 조합 폭발)", None, True),
    ("0", "아이디어: 비슷한 상황을 묶어 ‘상자(컨텍스트)’ 하나로 취급", GREEN, True),
    ("1", "예: ‘리버에서 / 좋은 패로 / 선공 위치 / 상대가 작게 레이즈한’ 상황 = 상자 1개"),
    ("0", "왜 이렇게까지 줄이나? — 신경망의 ‘근사 오차’를 없애려고", ACCENT, True),
    ("1", "신경망을 쓰면 ‘학습 실패’가 알고리즘 탓인지 환경 탓인지 구분 불가"),
    ("1", "상자를 256개로 고정하면 각 상자 방문 횟수를 ‘정확히’ 셀 수 있음"),
    ("0", "이렇게 단순화한 표(table) 위에서 학습 → ‘tabular’ 방식", None, True),
], top=Inches(1.45), gap=11, size=17)
s.notes_slide.notes_text_frame.text = (
    "추상화의 목적을 강조: 성능을 높이려는 게 아니라, 원인을 깨끗하게 분리하려는 통제 장치다."
)

# --- 배경 4: 256 상자를 만드는 4가지 기준 ---
s = add_slide()
title_bar(s, "0-4. 상자를 나누는 4가지 기준 (→ 256개)")
bullets(s, [
    ("0", "① 라운드 (4가지): 지금 몇 번째 베팅 단계인가", NAVY, True),
    ("1", "Preflop / Flop / Turn / River"),
    ("0", "② 핸드버킷 (8가지): 내 패가 얼마나 좋은가", NAVY, True),
    ("1", "PREMIUM(최상) → STRONG → … → POOR → TRASH(최악)"),
    ("0", "③ 위치 (2가지): 내가 선공인가 후공인가", NAVY, True),
    ("1", "BB(빅블라인드) / SB(스몰블라인드)"),
    ("0", "④ 직전 행동 (4가지): 상대가 방금 뭘 했나", NAVY, True),
    ("1", "없음 / 체크·콜 / 작은레이즈 / 큰레이즈"),
    ("0", "4 × 8 × 2 × 4 = 256개의 상황 상자(컨텍스트)", GREEN, True),
], top=Inches(1.3), gap=7, size=16)
s.notes_slide.notes_text_frame.text = (
    "네 가지 기준의 곱이 256. 뒤의 정밀 슬라이드(3. 상태 추상화)에서 행동 8을 곱해 2,048 셀이 된다."
)

# --- 배경 5: 고정 상대(페르소나) ---
s = add_slide()
title_bar(s, "0-5. 고정 상대(페르소나) — 학습하지 않는 ‘환경’")
bullets(s, [
    ("0", "상대는 학습하지 않음 → 정해진 스타일대로만 행동 (스크립트)", ACCENT, True),
    ("1", "그래서 상대는 ‘적’이 아니라 고정된 환경 규칙처럼 작동"),
    ("0", "5종 스타일(페르소나):", NAVY, True),
    ("1", "TAG: 타이트-공격 / LAG: 루즈-공격 / MAN: 매니악(과격)"),
    ("1", "STA: 콜링스테이션(잘 안 폴드) / NIT: 극단 보수(거의 폴드)"),
    ("0", "상대 스타일이 ‘어떤 상황이 생기는지’를 결정", GREEN, True),
    ("1", "예: NIT은 자주 폴드 → 특정 공격 상황 자체가 안 만들어짐"),
], top=Inches(1.5), gap=11, size=17)
s.notes_slide.notes_text_frame.text = (
    "‘상대=환경’ 개념의 구체화. 상대가 고정이라 P(s'|s,a)의 일부로 들어간다. "
    "이 5종이 만드는 전이 분포 차이가 뒤 실험의 독립변수."
)

# --- 배경 6: 도달 분포 d^π(s) ---
s = add_slide()
title_bar(s, "0-6. 도달 분포 d^π(s) — 각 상자에 얼마나 자주 가나")
bullets(s, [
    ("0", "d^π(s): 정책 π로 플레이할 때 상태(상자) s를 방문하는 빈도", NAVY, True),
    ("1", "쉽게: 256개 상자 각각에 ‘방문 도장’이 몇 번 찍히나"),
    ("0", "도장이 0번인 상자 = 한 번도 못 간 상황", RED, True),
    ("1", "표본이 없으니 학습 자체가 불가능 (배울 데이터가 없음)"),
    ("0", "우리는 신경망이 아니라 표라서 이 빈도를 ‘정확히’ 셀 수 있음", GREEN, True),
    ("1", "→ ‘못 간 상황’이 알고리즘 탓이 아니라 상대 탓임을 분리 증명 가능"),
], top=Inches(1.5), gap=12, size=18)
s.notes_slide.notes_text_frame.text = (
    "d^π(s)는 어려운 기호처럼 보이지만 ‘방문 빈도’일 뿐. 0번 방문 상자가 핵심 증거(45.5%)로 이어진다."
)

# --- 배경 7: MC vs Q-learning + Softmax + mbb/g ---
s = add_slide()
title_bar(s, "0-7. 학습 방식과 평가 단위 — 직관만")
bullets(s, [
    ("0", "몬테카를로(MC) 학습: 한 판을 ‘끝까지’ 가보고 최종 결과로 배움", NAVY, True),
    ("1", "Q-learning과 차이: 중간값 추정(부트스트랩) 없이 실제 결과만 사용"),
    ("1", "→ 추정 오차가 안 섞여서 ‘환경 탓’ 분리가 더 깨끗함"),
    ("0", "Softmax 탐색 + 온도(T): 얼마나 ‘모험적으로’ 행동할지 손잡이", NAVY, True),
    ("1", "온도 높음=다양하게 시도 / 낮음=잘하던 것 위주 (10.0→0.5로 감소)"),
    ("0", "mbb/g: 성능 단위 (한 판당 빅블라인드의 1/1000을 얼마나 벌었나)", NAVY, True),
    ("1", "+면 이득, −면 손해. 예: +1083.9면 크게 이득, −650.4면 크게 손해"),
], top=Inches(1.4), gap=9, size=16)
s.notes_slide.notes_text_frame.text = (
    "세 개념을 한 장에 압축. ‘온도’가 곧 탐색 강도 손잡이라는 점을 강조 → 뒤 6-B 실험과 직결. "
    "mbb/g 부호(+/−)만 읽으면 결과 표를 이해할 수 있다고 안내."
)

# ---------------------------------------------------------------------------
# Slide 4 — 설계 아이디어: 근사 오차 배제
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "2. 설계 아이디어 — 근사 오차를 제거한 통제 환경",
          "(앞의 0-3 직관을 정식으로 정리)")
bullets(s, [
    ("0", "심층 신경망은 ‘일반화 실패’와 ‘환경 차단’을 분리할 수 없음", None, True),
    ("0", "해결: HUNL을 순수 표-기반(tabular) 환경으로 극단 추상화", GREEN, True),
    ("1", "근사 오차가 개입할 수 없는 통제 조건 확보"),
    ("1", "상태 도달 분포 d^π(s)를 셀 단위로 전수(exact) 측정 가능"),
    ("0", "이로써 커버리지 한계의 원인을 깨끗하게 분리 관측", None, True),
    ("1", "= 고정 상대가 만드는 외생적 전이 분포에 병목이 있는가?"),
], top=Inches(1.5), gap=12)
s.notes_slide.notes_text_frame.text = (
    "핵심 방법론적 강점: 256셀 exact 측정으로 도달 실패의 원인을 수리적으로 분리한다."
)

# ---------------------------------------------------------------------------
# Slide 5 — 추상화 구조도
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "3. 상태 추상화 — 256 컨텍스트 × 8 행동 = 2,048 셀")


def factor_box(slide, left, top, w, h, title, lines, fill):
    b = slide.shapes.add_shape(1, left, top, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill; b.line.color.rgb = NAVY
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; set_font(r, 15, True, NAVY)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = lines; set_font(r2, 11, False, GRAY)


top = Inches(1.6); h = Inches(1.2); w = Inches(2.7)
factor_box(s, Inches(0.6), top, w, h, "라운드 (4)", "Preflop/Flop/Turn/River", LIGHT)
factor_box(s, Inches(3.5), top, w, h, "핸드버킷 (8)", "PREMIUM…TRASH", LIGHT)
factor_box(s, Inches(6.4), top, w, h, "위치 (2)", "BB / SB", LIGHT)
factor_box(s, Inches(9.3), top, w, h, "직전행동 (4)", "NONE/CC/SR/BR", LIGHT)
# 곱 표시
tf = add_textbox(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.6))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "4 × 8 × 2 × 4 = 256 컨텍스트"
set_font(r, 20, True, ACCENT)
factor_box(s, Inches(3.5), Inches(3.8), Inches(2.7), Inches(1.1),
           "행동 (8)", "FOLD/CHECK/CALL/R25/50/75/100/ALLIN", RGBColor(0xE7, 0xF1, 0xE7))
tf = add_textbox(s, Inches(6.4), Inches(3.8), Inches(6), Inches(1.1), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "256 × 8 = 2,048 셀  → d^π(s) 전수 측정 대상"
set_font(r, 18, True, GREEN)
bullets(s, [
    ("0", "도메인: HUNL, 시작스택 200, SB=1/BB=2, pokerkit 엔진"),
    ("0", "핸드버킷은 256 컨텍스트 안에 이미 포함 (별도 추가 아님)"),
], top=Inches(5.3), size=15, gap=8)
s.notes_slide.notes_text_frame.text = (
    "여기서 cardinality를 정확히: 256 컨텍스트 안에 핸드버킷 8이 인수로 들어가 있고, "
    "행동 8을 곱해 총 2,048 셀이다."
)

# ---------------------------------------------------------------------------
# Slide 6 — 알고리즘
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "3. 학습 알고리즘 — On-policy Monte Carlo (Q-learning 아님)")
bullets(s, [
    ("0", "Q-table 자료구조를 on-policy Monte Carlo로 갱신", None, True),
    ("1", "갱신식:  Q(s,a) ← Q(s,a) + α · ( G − Q(s,a) )   (부트스트랩 없음)"),
    ("1", "에피소드 실제 리턴 G로 갱신 → off-policy Q-learning 부트스트랩 미사용"),
    ("0", "탐색: Softmax (온도 10.0 → 0.5, 학습 전반 80% 구간 선형 감쇠)", None, True),
    ("1", "PrevAction 특징 포함, CHECK=1chip 처리"),
    ("0", "하이퍼파라미터: α=0.1, γ=0.9, seed=42", None, True),
    ("0", "평가: 어댑터 끈 raw-greedy, 100k 핸드 × 다중 시드, 학습/평가 분리", None, True),
], top=Inches(1.5), gap=11)
s.notes_slide.notes_text_frame.text = (
    "질문 대비: 코드에 update_q(Q-learning)도 있으나 미사용. 실제 사용은 update_mc. "
    "정확한 명칭은 'tabular on-policy Monte Carlo control with softmax exploration'."
)

# ---------------------------------------------------------------------------
# Slide 7 — 처방 3종 기각 (개요)
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "4. 통념적 처방 3종을 차례로 적용 → 모두 기각")
bullets(s, [
    ("0", "미방문 셀이 ‘과적합/학습부족’ 때문이라는 의심을 검증", None, True),
    ("0", "[6-A] 예산 부족 가설 → 기각", RED, True),
    ("1", "도달 확률이 정확히 0인 셀이 전체의 45.5% (T-UNREACH, 버그 아님)"),
    ("1", "예산을 무한히 늘려도 리턴 표본이 없어 갱신 0회 → 안 채워짐"),
    ("0", "[6-B] 탐색(온도) 강화 가설 → 기각", RED, True),
    ("1", "온도 4방향 스윕에도 미학습 컨텍스트 67 → 68/71/72 (거의 불변)"),
    ("0", "[6-C] ‘넓은 커버리지=흑자’ 가설 → 기각", RED, True),
    ("1", "흑자 run이 baseline보다 미학습 셀이 오히려 더 많은 경우 관측"),
], top=Inches(1.4), gap=9, size=17)
s.notes_slide.notes_text_frame.text = (
    "세 처방 모두 천장을 못 움직였다. 병목은 탐색 강도가 아니라 상대가 만드는 전이 분포."
)

# ---------------------------------------------------------------------------
# Slide 8 — 처방 기각 수치 강조
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "4. 핵심 수치 — 탐색을 흔들어도 천장은 불변")


def stat_card(slide, left, top, big, small, color):
    b = slide.shapes.add_shape(1, left, top, Inches(3.7), Inches(2.4))
    b.fill.solid(); b.fill.fore_color.rgb = LIGHT; b.line.color.rgb = color
    b.line.width = Pt(2)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big; set_font(r, 32, True, color)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = small; set_font(r2, 14, False, GRAY)


stat_card(s, Inches(0.6), Inches(2.0), "45.5%", "도달 확률 0인 셀 비율\n(예산으로 못 채움)", RED)
stat_card(s, Inches(4.8), Inches(2.0), "67 → 72", "온도 4방향 스윕 후\n미학습 컨텍스트(거의 불변)", ACCENT)
stat_card(s, Inches(9.0), Inches(2.0), "흑자 ≠ 넓이", "흑자 원천은 커버리지 넓이\n아니라 상대 분포 정합성", GREEN)
bullets(s, [
    ("0", "탐색 온도는 ‘이미 도달한 셀 내부’의 행동 다양성만 키움", None, True),
    ("0", "셀 자체의 도달 확률은 바꾸지 못함 → 천장은 환경이 결정", ACCENT, True),
], top=Inches(4.9), size=18, gap=10)
s.notes_slide.notes_text_frame.text = (
    "이 세 숫자가 처방 3종 기각의 핵심 증거. 수치는 절대 바꾸지 말 것."
)

# ---------------------------------------------------------------------------
# Slide 9 — 반대 증거: 분포 다양화 (표)
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "4. 반대 증거 — 상대 분포만 바꾸니 천장이 움직였다",
          "알고리즘 수정 없이 ‘학습 상대 분포’만 다양화 (각 2M 에피소드, 100k×5 평가)")
rows = [
    ["Run", "vs Random (mbb/g, SD)", "vs 고정 TAG (mbb/g, SD)", "판정"],
    ["26a LAG 단일 (참고)", "−650.4 (19.3)", "+924.4 (12.2)", "vs Random 적자"],
    ["26c STA 단일 (참고)", "−674.9 (6.5)", "+774.3 (1.4)", "vs Random 적자"],
    ["27a CYCLE (순환)", "+681.3 (51.4)", "+868.3 (12.1)", "동시 흑자"],
    ["27b MIXED (확률혼합)", "+1083.9 (110.5)", "+885.2 (10.9)", "동시 흑자"],
]
table = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.9),
                           Inches(12.1), Inches(3.4)).table
table.columns[0].width = Inches(3.2)
table.columns[1].width = Inches(3.3)
table.columns[2].width = Inches(3.3)
table.columns[3].width = Inches(2.3)
for ci in range(4):
    cell = table.cell(0, ci)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = rows[0][ci]
    set_font(r, 13, True, RGBColor(0xFF, 0xFF, 0xFF))
for ri in range(1, len(rows)):
    highlight = ri >= 3
    for ci in range(4):
        cell = table.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE7, 0xF1, 0xE7) if highlight else RGBColor(0xFF, 0xFF, 0xFF)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = rows[ri][ci]
        col = GREEN if highlight else GRAY
        set_font(r, 12, highlight, col)
bullets(s, [
    ("0", "적자였던 LAG·STA를 학습 분포에 섞으니 두 상대 모두 동시 흑자 → 분포 강건성", GREEN, True),
    ("0", "27b vs Random 하한 1083.9 − 3×110.5 = 752.4 > 0 → 부호 안정성 확보", None, True),
], top=Inches(5.5), size=15, gap=8)
s.notes_slide.notes_text_frame.text = (
    "처방으로 안 움직이던 천장이 ‘상대 분포’ 레버로 움직였다는 결정적 대조. 표의 모든 수치 고정."
)

# ---------------------------------------------------------------------------
# Slide 10 — 논의 / 정직한 한계
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "5. 논의 — 무엇을 주장하고, 무엇을 주장하지 않는가")
bullets(s, [
    ("0", "주장하는 것", GREEN, True),
    ("1", "시험한 레버(예산·탐색·커버리지) 중 ‘상대 전이 분포’만 천장을 움직였다"),
    ("1", "커버리지 통제권이 학습자 → 환경으로 넘어가는 경계 조건을 실증"),
    ("1", "기존 집중성/커버리지 이론과 ‘일치’ — 숨은 전제(통제권=학습자)를 명료화"),
    ("0", "주장하지 않는 것 (정직한 한계)", RED, True),
    ("1", "‘전이 분포가 천장을 유일하게 결정한다’는 강한 명제 → 범위 밖"),
    ("1", "본 결과는 tabular · on-policy MC 한정 조건부 명제"),
    ("0", "기여 성격: 경계 조건 명료화 + 실증적·해석적 신규성", ACCENT, True),
], top=Inches(1.45), gap=9, size=16)
s.notes_slide.notes_text_frame.text = (
    "이 슬라이드가 honesty의 핵심. 결론에서도 이 hedge를 절대 깨지 않는다."
)

# ---------------------------------------------------------------------------
# Slide 11 — 결론 & 향후 과제
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "6. 결론 및 향후 과제")
bullets(s, [
    ("0", "결론", NAVY, True),
    ("1", "256셀 전수 측정으로 ‘탐색 강화 = 커버리지 확장’ 통념을 조건부 기각"),
    ("1", "고정 상대가 외생적으로 커버리지 천장을 결정함을 실증"),
    ("1", "흑자의 원천은 커버리지 넓이가 아니라 상대 분포 정합성"),
    ("0", "향후 과제", ACCENT, True),
    ("1", "버킷 세분화(8→16) 후에도 빈 셀이 상대 종속으로 남는지 재확인"),
    ("1", "함수 근사 / off-policy(CFR 계열) / 리그 self-play 등 미검증 레버 시험"),
    ("1", "상대도 학습하는 다중 에이전트 환경에서 통제권 이동 재현 여부"),
], top=Inches(1.45), gap=10, size=17)
s.notes_slide.notes_text_frame.text = (
    "결론도 '조건부 기각'으로 못박는다. 다중 에이전트는 향후 일반화 함의로만 언급."
)

# ---------------------------------------------------------------------------
# Slide 12 — Q&A 대비 (백업)
# ---------------------------------------------------------------------------
s = add_slide()
title_bar(s, "예상 질문 (Backup)")
bullets(s, [
    ("0", "Q. 왜 Q-learning이 아니라 Monte Carlo인가?", NAVY, True),
    ("1", "A. 부트스트랩 오차 배제, 실제 리턴 G로 갱신해 도달 분리를 깨끗하게."),
    ("0", "Q. 256셀은 과한 추상화 아닌가?", NAVY, True),
    ("1", "A. 근사 오차를 배제해 환경 차단을 분리하려는 의도적 통제. 세분화는 향후 과제."),
    ("0", "Q. 다중 에이전트면 결과가 달라지지 않나?", NAVY, True),
    ("1", "A. 본 명제는 단일 에이전트 MDP 한정. 다중 에이전트는 향후 검증 대상으로 명시."),
    ("0", "Q. ‘유일성’은 어떻게 보장하나?", NAVY, True),
    ("1", "A. 보장하지 않음. 강한 명제는 범위 밖이며 조건부 명제로만 주장."),
], top=Inches(1.45), gap=9, size=16)
s.notes_slide.notes_text_frame.text = "발표 시간 부족하면 생략, 질문 들어오면 표시."

import os
out = os.path.join(os.path.dirname(__file__), "발표_커버리지천장.pptx")
try:
    prs.save(out)
except PermissionError:
    import datetime
    stamp = datetime.datetime.now().strftime("%H%M%S")
    out = os.path.join(os.path.dirname(__file__), f"발표_커버리지천장_{stamp}.pptx")
    prs.save(out)
    print("원본이 열려 있어 새 파일로 저장했습니다.")
print("saved:", out)
print("slides:", len(prs.slides._sldIdLst))

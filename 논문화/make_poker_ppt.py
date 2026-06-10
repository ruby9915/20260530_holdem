# -*- coding: utf-8 -*-
"""
make_poker_ppt.py
27번 실험(혼합 페르소나 RL 포커) 발표용 PPTX 생성.
예시(Sal_넥서스 로봇팔) 양식: 16:9, 맑은 고딕, 네이비/블루, 섹션 헤더 + 푸터.
이미지 대신 matplotlib 그래프/도식 사용.
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

matplotlib.rcParams["font.family"] = "Malgun Gothic"
matplotlib.rcParams["axes.unicode_minus"] = False

ASSET = os.path.join(os.path.dirname(__file__), "_ppt_assets")
os.makedirs(ASSET, exist_ok=True)

# ── 팔레트 ────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
LBLUE  = RGBColor(0xE9, 0xF1, 0xFA)
GRAY   = RGBColor(0x59, 0x59, 0x59)
DARK   = RGBColor(0x22, 0x22, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xE8, 0xA1, 0x3A)

C_NAVY  = "#1F3864"
C_BLUE  = "#2E6DB4"
C_LBLUE = "#A9C7E8"
C_ACC   = "#E8A13A"
C_GREEN = "#3FA34D"
C_RED   = "#C0392B"
C_GRAY  = "#8895A7"

# ═══════════════════════════════════════════════════════
# 1. 그래프/도식 생성
# ═══════════════════════════════════════════════════════
def fig_concept():
    """단일 상대 학습 vs 혼합 상대 학습 대비 도식."""
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.6))
    for ax in axes:
        ax.set_xlim(0, 10); ax.set_ylim(0, 10); ax.axis("off")

    # (좌) 단일 상대
    ax = axes[0]
    ax.set_title("단일 상대 학습  →  과적합", fontsize=15, color=C_RED, fontweight="bold", pad=12)
    ag = FancyBboxPatch((4, 4.2), 2, 1.6, boxstyle="round,pad=0.1",
                        fc=C_NAVY, ec="none"); ax.add_patch(ag)
    ax.text(5, 5, "에이전트", color="white", ha="center", va="center", fontsize=12, fontweight="bold")
    opp = FancyBboxPatch((4, 8), 2, 1.2, boxstyle="round,pad=0.1", fc=C_GRAY, ec="none")
    ax.add_patch(opp)
    ax.text(5, 8.6, "TAG 한 명", color="white", ha="center", va="center", fontsize=11)
    ax.add_patch(FancyArrowPatch((5, 7.9), (5, 5.9), arrowstyle="-|>", mutation_scale=20, lw=2, color=C_RED))
    ax.text(7.4, 6.9, "좁은 경험", color=C_RED, fontsize=11, ha="center")
    ax.text(5, 2.7, "처음 보는 상대에겐 약함", color=C_RED, ha="center", fontsize=11, style="italic")

    # (우) 혼합 상대
    ax = axes[1]
    ax.set_title("혼합 상대 학습  →  분포 강건", fontsize=15, color=C_GREEN, fontweight="bold", pad=12)
    ag = FancyBboxPatch((4, 4.2), 2, 1.6, boxstyle="round,pad=0.1", fc=C_NAVY, ec="none")
    ax.add_patch(ag)
    ax.text(5, 5, "에이전트", color="white", ha="center", va="center", fontsize=12, fontweight="bold")
    names = ["TAG", "LAG", "MAN", "STA", "NIT"]
    cols  = [C_BLUE, C_ACC, C_RED, C_GREEN, C_GRAY]
    pos = [(1.2, 8.2), (3.4, 9.0), (5.6, 9.0), (7.8, 8.2), (8.6, 5.2)]
    for (x, y), nm, c in zip(pos, names, cols):
        b = FancyBboxPatch((x-0.7, y-0.5), 1.4, 1.0, boxstyle="round,pad=0.08", fc=c, ec="none")
        ax.add_patch(b)
        ax.text(x, y, nm, color="white", ha="center", va="center", fontsize=10, fontweight="bold")
        ax.add_patch(FancyArrowPatch((x, y-0.55), (5, 5.9), arrowstyle="-|>",
                                     mutation_scale=13, lw=1.5, color=c, alpha=0.8))
    ax.text(5, 2.7, "누구에게나 통하는 정책", color=C_GREEN, ha="center", fontsize=11, style="italic")

    plt.tight_layout()
    p = os.path.join(ASSET, "concept.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_state_action():
    """상태 4축(=256) × 행동 8 = 2,048칸 Q-테이블 분해 도식."""
    fig, ax = plt.subplots(figsize=(11, 4.4))
    ax.set_xlim(0, 24); ax.set_ylim(0, 10); ax.axis("off")

    # 상태 4개 축 박스
    dims = [("Round", 4), ("Position", 2), ("핸드강도", 8), ("PrevAction", 4)]
    cols = [C_BLUE, C_GREEN, C_ACC, C_NAVY]
    x = 0.4; w = 2.7; gap = 0.5
    centers = []
    for (nm, v), c in zip(dims, cols):
        b = FancyBboxPatch((x, 5.4), w, 2.6, boxstyle="round,pad=0.1", fc=c, ec="none")
        ax.add_patch(b)
        ax.text(x+w/2, 7.0, nm, color="white", ha="center", va="center", fontsize=12, fontweight="bold")
        ax.text(x+w/2, 6.1, f"{v}", color="white", ha="center", va="center", fontsize=15, fontweight="bold")
        centers.append(x+w)
        if c is not cols[-1]:
            ax.text(x+w+gap/2, 6.7, "×", ha="center", va="center", fontsize=18, color=C_GRAY, fontweight="bold")
        x += w + gap
    ax.text(x+0.2, 6.7, "=", ha="center", va="center", fontsize=18, color=C_NAVY, fontweight="bold")
    sb = FancyBboxPatch((x+0.6, 5.4), 3.0, 2.6, boxstyle="round,pad=0.1", fc="#34507a", ec="none")
    ax.add_patch(sb)
    ax.text(x+0.6+1.5, 7.0, "상태 수", color="white", ha="center", va="center", fontsize=12, fontweight="bold")
    ax.text(x+0.6+1.5, 6.0, "256", color="white", ha="center", va="center", fontsize=20, fontweight="bold")

    # 아래: 256 상태 × 8 행동 = 2,048 Q-테이블
    ax.text(12, 3.7, "256 상태", ha="center", va="center", fontsize=14, color=C_NAVY, fontweight="bold")
    ax.text(13.6, 3.7, "×", ha="center", va="center", fontsize=16, color=C_GRAY, fontweight="bold")
    ax.text(15.2, 3.7, "행동 8종", ha="center", va="center", fontsize=14, color=C_ACC, fontweight="bold")
    ax.text(17.0, 3.7, "=", ha="center", va="center", fontsize=16, color=C_NAVY, fontweight="bold")
    qb = FancyBboxPatch((17.6, 2.9), 5.6, 1.7, boxstyle="round,pad=0.12", fc=C_GREEN, ec="none")
    ax.add_patch(qb)
    ax.text(20.4, 3.75, "Q-테이블  2,048칸", color="white", ha="center", va="center", fontsize=15, fontweight="bold")
    ax.text(12, 1.4, "각 (상태, 행동) 칸마다 가치 Q(s, a) 하나를 저장 → 표 한 장으로 정책 전체 표현",
            ha="center", va="center", fontsize=11.5, color=C_GRAY, style="italic")
    ax.set_title("상태 256개 × 행동 8개  →  2,048칸 Q-테이블",
                 fontsize=15, color=C_NAVY, fontweight="bold", pad=6)
    p = os.path.join(ASSET, "state_action.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_reward_split():
    """비례 배분 보상: 구체 예시(투입 칩 비율로 payoff 분배)."""
    fig, ax = plt.subplots(figsize=(11, 4.6))
    ax.set_xlim(0, 24); ax.set_ylim(0, 10); ax.axis("off")

    # 한 핸드의 세 번의 결정과 투입 칩
    steps = [("PREFLOP\n레이즈", 10), ("FLOP\n베팅", 20), ("TURN\n큰 베팅", 60)]
    total = sum(v for _, v in steps)       # 90
    payoff = 180                            # 최종 순이익 +180
    mult = payoff / total                   # 2.0
    cols = [C_BLUE, C_ACC, C_RED]

    ax.text(5.0, 9.3, "한 핸드에서 내린 3번의 결정 (투입 칩)", ha="center", fontsize=13,
            color=C_NAVY, fontweight="bold")
    x = 0.6; w = 2.9; gap = 0.6; centers = []
    for (nm, v), c in zip(steps, cols):
        h = 0.9 + v/60*2.6
        b = FancyBboxPatch((x, 5.2), w, h, boxstyle="round,pad=0.08", fc=c, ec="none")
        ax.add_patch(b)
        ax.text(x+w/2, 5.2+h/2+0.15, nm, color="white", ha="center", va="center", fontsize=11, fontweight="bold")
        ax.text(x+w/2, 5.2+h/2-0.55, f"{v}칩", color="white", ha="center", va="center", fontsize=11)
        centers.append((x+w/2, c, v, nm))
        x += w + gap
    ax.text(x+0.1, 6.3, f"합계\n{total}칩", ha="center", va="center", fontsize=12, color=C_NAVY, fontweight="bold")

    # 핸드 결과
    res = FancyBboxPatch((12.6, 7.4), 10.8, 1.6, boxstyle="round,pad=0.1", fc="#2b6b3a", ec="none")
    ax.add_patch(res)
    ax.text(18.0, 8.2, f"핸드 종료 → 순이익  payoff = +{payoff}칩  (승리)",
            color="white", ha="center", va="center", fontsize=13, fontweight="bold")

    # 배분 결과 화살표 + 각 스텝 받는 보상
    ax.text(18.0, 6.4, "투입한 칩 비율 그대로 +180을 나눠 가짐", ha="center", fontsize=12,
            color=C_NAVY, fontweight="bold")
    bx = 12.8; bw = 3.2; bgap = 0.5
    for (cx, c, v, nm), i in zip(centers, range(3)):
        g = v/total*payoff
        b = FancyBboxPatch((bx, 3.4), bw, 1.7, boxstyle="round,pad=0.08", fc=c, ec="none")
        ax.add_patch(b)
        ax.text(bx+bw/2, 4.55, nm.split(chr(10))[0], color="white", ha="center", va="center", fontsize=10, fontweight="bold")
        ax.text(bx+bw/2, 3.85, f"+{g:.0f}", color="white", ha="center", va="center", fontsize=15, fontweight="bold")
        ax.add_patch(FancyArrowPatch((cx, 5.1), (bx+bw/2, 5.15), arrowstyle="-|>",
                     mutation_scale=14, lw=1.6, color=c, alpha=0.7))
        bx += bw + bgap

    ax.text(12.0, 1.9, "큰 베팅(60칩)을 한 TURN 결정이 +120으로 가장 큰 책임을 받는다",
            ha="center", fontsize=12, color=C_RED, fontweight="bold")
    ax.text(12.0, 1.0, "→ 결과를 만든 '핵심 결정'이 더 크게 학습된다", ha="center",
            fontsize=11.5, color=C_GRAY, style="italic")
    p = os.path.join(ASSET, "reward_split.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_qflow():
    """Q-update(MC) 흐름도."""
    fig, ax = plt.subplots(figsize=(11, 3.2))
    ax.set_xlim(0, 22); ax.set_ylim(0, 6); ax.axis("off")
    steps = ["1핸드\n플레이", "각 스텝\n(s, a, invest)\n기록", "핸드 종료\npayoff 확정",
             "투입 비율로\nreturn g 배분", "모든 스텝\nQ(s,a)\n업데이트"]
    cols = [C_BLUE, C_BLUE, C_ACC, C_NAVY, C_GREEN]
    x = 0.5; w = 3.6; gap = 0.7
    centers = []
    for s, c in zip(steps, cols):
        b = FancyBboxPatch((x, 1.6), w, 2.8, boxstyle="round,pad=0.1", fc=c, ec="none")
        ax.add_patch(b)
        ax.text(x+w/2, 3.0, s, color="white", ha="center", va="center", fontsize=11, fontweight="bold")
        centers.append(x+w)
        x += w + gap
    for cx in centers[:-1]:
        ax.add_patch(FancyArrowPatch((cx+0.04, 3.0), (cx+gap-0.04, 3.0),
                     arrowstyle="-|>", mutation_scale=18, lw=2, color=C_GRAY))
    plt.tight_layout()
    p = os.path.join(ASSET, "qflow.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_temp():
    """온도 스케줄 + softmax 개념."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.9))
    # 온도 스케줄
    ep = np.linspace(0, 1, 200)
    tau = np.where(ep < 0.8, 10.0 + (0.5-10.0)*(ep/0.8), 0.5)
    ax1.plot(ep*100, tau, color=C_BLUE, lw=3)
    ax1.fill_between(ep*100, tau, color=C_BLUE, alpha=0.12)
    ax1.set_xlabel("학습 진행도 (%)", fontsize=11)
    ax1.set_ylabel("온도 τ", fontsize=11)
    ax1.set_title("온도 스케줄: 탐색→활용", fontsize=13, color=C_NAVY, fontweight="bold")
    ax1.annotate("초반 τ=10\n(탐색)", (5, 9.3), fontsize=10, color=C_RED)
    ax1.annotate("후반 τ=0.5\n(활용)", (60, 1.6), fontsize=10, color=C_GREEN)
    ax1.spines[["top", "right"]].set_visible(False)
    # softmax 확률 예시
    acts = ["FOLD", "CHECK", "CALL", "R50", "R100"]
    q = np.array([-2, 1, 3, 5, 2])
    for tau_v, c, lbl in [(5.0, C_ACC, "τ=5 (탐색적)"), (0.5, C_BLUE, "τ=0.5 (탐욕적)")]:
        e = np.exp((q-q.max())/tau_v); pr = e/e.sum()
        ax2.plot(acts, pr, "o-", color=c, lw=2.5, ms=7, label=lbl)
    ax2.set_ylabel("선택 확률 P(a)", fontsize=11)
    ax2.set_title("Softmax 행동 선택", fontsize=13, color=C_NAVY, fontweight="bold")
    ax2.legend(fontsize=10); ax2.spines[["top", "right"]].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, "temp.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_results():
    """결과 그래프: cycle/mixed vs Random/TAG."""
    fig, ax = plt.subplots(figsize=(9.8, 4.6))
    groups = ["cycle", "mixed"]
    rand = [681, 1084]; tag = [868, 885]
    x = np.arange(len(groups)); w = 0.35
    b1 = ax.bar(x-w/2, rand, w, label="vs Random (미학습 상대)", color=C_BLUE)
    b2 = ax.bar(x+w/2, tag, w, label="vs 고정 TAG (정석 상대)", color=C_ACC)
    for bars in (b1, b2):
        for b in bars:
            ax.text(b.get_x()+b.get_width()/2, b.get_height()+15,
                    f"+{int(b.get_height())}", ha="center", fontsize=12, fontweight="bold")
    ax.axhline(0, color="k", lw=1)
    ax.axhspan(-80, 0, color=C_RED, alpha=0.06)
    ax.set_xticks(x); ax.set_xticklabels(groups, fontsize=13, fontweight="bold")
    ax.set_ylabel("수익률 mbb/g", fontsize=12)
    ax.set_ylim(0, 1250)
    ax.set_title("두 방식 모두 양쪽 상대에서 동시 흑자 → 분포 강건성 달성",
                 fontsize=14, color=C_NAVY, fontweight="bold", pad=12)
    ax.legend(fontsize=11, loc="upper left")
    ax.spines[["top", "right"]].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, "results.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


def fig_personas():
    """5종 페르소나 공격성 스펙트럼."""
    fig, ax = plt.subplots(figsize=(10, 3.6))
    names = ["STA\nCalling Station", "NIT\nRock", "TAG\nTight-Aggr", "LAG\nLoose-Aggr", "MAN\nManiac"]
    aggr  = [5, 12, 21, 30, 40]   # 레이즈 비중(%) 근사
    cols  = [C_GREEN, C_GRAY, C_BLUE, C_ACC, C_RED]
    bars = ax.barh(names, aggr, color=cols)
    for b, v in zip(bars, aggr):
        ax.text(v+0.6, b.get_y()+b.get_height()/2, f"{v}%", va="center", fontsize=11, fontweight="bold")
    ax.set_xlabel("공격성 (레이즈 비중, 근사)", fontsize=11)
    ax.set_xlim(0, 46)
    ax.set_title("학습 상대 5종 — 서로 다른 전이분포 제공", fontsize=14, color=C_NAVY, fontweight="bold", pad=10)
    ax.spines[["top", "right"]].set_visible(False)
    plt.tight_layout()
    p = os.path.join(ASSET, "personas.png"); fig.savefig(p, dpi=160, bbox_inches="tight"); plt.close(fig)
    return p


# ═══════════════════════════════════════════════════════
# 2. PPTX 빌드
# ═══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Emu(6858000)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
FONT = "맑은 고딕"


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def add_textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def footer(slide):
    tf = add_textbox(slide, Inches(0.45), SH-Inches(0.5), Inches(6), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Dongseo University Soft  ·  Reinforcement Learning"
    _set(r, 10, GRAY)


def content_header(slide, section, title):
    """예시 양식: 좌상단 섹션 + 큰 제목 + 상단 액센트 바."""
    rect(slide, 0, 0, SW, Inches(1.5), WHITE)
    rect(slide, 0, Inches(1.5), SW, Pt(3), BLUE)
    rect(slide, Inches(0.45), Inches(0.35), Pt(6), Inches(0.85), ACCENT)
    tf = add_textbox(slide, Inches(0.7), Inches(0.28), Inches(11), Inches(0.5))
    r = tf.paragraphs[0].add_run(); r.text = section; _set(r, 14, BLUE, bold=True)
    tf2 = add_textbox(slide, Inches(0.7), Inches(0.68), Inches(11), Inches(0.7))
    r = tf2.paragraphs[0].add_run(); r.text = title; _set(r, 26, NAVY, bold=True)
    footer(slide)


def bullets(slide, items, l, t, w, h, size=16, gap=10):
    tf = add_textbox(slide, l, t, w, h)
    first = True
    for it in items:
        lvl = it[0] if isinstance(it, tuple) else 0
        txt = it[1] if isinstance(it, tuple) else it
        bold = it[2] if (isinstance(it, tuple) and len(it) > 2) else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        mark = "■  " if lvl == 0 else "•  "
        r = p.add_run(); r.text = mark + txt
        _set(r, size if lvl == 0 else size-2, NAVY if lvl == 0 else DARK, bold=(lvl == 0 or bold))
    return tf


def pic_fit(slide, path, l, t, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw/ih; bw, bh = max_w, int(max_w/ar)
    if bh > max_h:
        bh = max_h; bw = int(max_h*ar)
    left = l + (max_w-bw)//2
    slide.shapes.add_picture(path, left, t, bw, bh)


# ── 슬라이드 1: 타이틀 ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.3), SW, Pt(4), ACCENT)
tf = add_textbox(s, Inches(1.0), Inches(1.7), Inches(10.2), Inches(2.4))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "상대를 가리지 않는 포커 AI"; _set(r, 44, WHITE, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(14); r = p2.add_run()
r.text = "학습 상대 다양화를 통한 분포 강건 강화학습"; _set(r, 24, RGBColor(0xBD, 0xD3, 0xEE))
tf2 = add_textbox(s, Inches(1.0), Inches(4.7), Inches(10), Inches(1.2))
p = tf2.paragraphs[0]; r = p.add_run()
r.text = "헤즈업 No-Limit 텍사스 홀덤  ·  Monte-Carlo Q-Learning"; _set(r, 16, RGBColor(0xD7, 0xE3, 0xF4))
p = tf2.add_paragraph(); p.space_before = Pt(8); r = p.add_run()
r.text = "강화학습 기말 발표"; _set(r, 14, RGBColor(0xA9, 0xC0, 0xE0))

# ── 슬라이드 2: 배경 ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "배경", "포커는 왜 어려운 강화학습 문제인가")
bullets(s, [
    "불완전정보 게임 : 상대의 카드를 볼 수 없다 → 체스·바둑(완전정보)과 근본적으로 다름",
    (1, "내가 가진 정보만으로 최선의 베팅을 추론해야 함"),
    "확률적·고분산 보상 : 같은 선택도 카드 운에 따라 결과가 크게 흔들림",
    (1, "한 판 결과만으로는 그 선택이 좋았는지 알기 어려움"),
    "상대 의존성 : 최적 전략이 상대 성향에 따라 달라짐",
    (1, "공격적 상대 ↔ 소극적 상대 → 정답 액션이 정반대가 되기도 함"),
], Inches(0.7), Inches(1.85), Inches(7.4), Inches(4.3), size=16, gap=8)
# 우측 환경 박스
rect(s, Inches(8.5), Inches(2.0), Inches(3.1), Inches(3.4), LBLUE)
tf = add_textbox(s, Inches(8.7), Inches(2.15), Inches(2.7), Inches(3.1))
p = tf.paragraphs[0]; r = p.add_run(); r.text = "실험 환경"; _set(r, 16, NAVY, bold=True)
for line in ["• 1:1 헤즈업 홀덤", "• No-Limit", "• 시작 스택 200칩", "• 빅블라인드 2칩", "  (= 100BB 깊이)", "• 1 핸드 = 1 에피소드"]:
    p = tf.add_paragraph(); p.space_before = Pt(7); r = p.add_run(); r.text = line; _set(r, 13, DARK)

# ── 슬라이드 3: 제안 ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "제안 내용", "학습 상대를 다양화해 '분포 강건성'을 학습")
bullets(s, [
    "문제의식 : 한 종류 상대로만 훈련하면 그 상대에게만 강하다 (분포 과적합)",
    "제안 : 매 에피소드마다 상대 성향(persona)을 바꿔가며 훈련",
    (1, "여러 상황 분포를 두루 경험 → 누구에게나 통하는 정책"),
], Inches(0.7), Inches(1.8), Inches(11), Inches(1.7), size=16, gap=8)
pic_fit(s, fig_concept(), Inches(0.7), Inches(3.4), Inches(11), Inches(3.0))

# ── 슬라이드 4: 모델링 ① State & Action ──────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "모델링 ①  ·  State & Action", "상태 256개 × 행동 8개 = 2,048칸 Q-테이블")
pic_fit(s, fig_state_action(), Inches(0.6), Inches(1.7), Inches(11), Inches(3.3))
bullets(s, [
    "상태(State) = 4개 축의 조합 : Round(4) × Position(2) × 핸드강도(8) × PrevAction(4) = 256",
    (1, "PrevAction(직전 상대 행동) 축으로 '상대가 크게 베팅했는가'에 조건부 반응"),
    "행동(Action) 8종 : FOLD / CHECK / CALL / RAISE 25·50·75·100% / ALL-IN",
    (1, "매 턴 합법적인 행동만 후보 → 256 상태마다 8행동의 가치 Q(s,a)를 표로 저장"),
], Inches(0.7), Inches(5.05), Inches(11), Inches(1.6), size=14, gap=6)

# ── 슬라이드 5: 모델링 ② Reward 배분 ─────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "모델링 ②  ·  Reward (보상 배분)", "한 핸드의 순이익을 '투입한 칩 비율'로 나눠 준다")
bullets(s, [
    "포커는 한 핸드가 끝나야 손익이 정해진다 → 보상은 핸드 종료 시 한 번, '딴 칩 − 잃은 칩'으로 발생",
    "그런데 한 핸드 동안 여러 번 결정을 내린다 → 이 하나의 결과를 각 결정에 어떻게 나눠줄까?",
    "해법 : 각 결정이 베팅에 넣은 칩의 비율만큼 결과를 나눠 갖는다 (많이 건 결정이 더 큰 책임)",
], Inches(0.7), Inches(1.75), Inches(11), Inches(1.9), size=15, gap=8)
pic_fit(s, fig_reward_split(), Inches(0.6), Inches(3.55), Inches(11), Inches(2.85))

# ── 슬라이드 6: 모델링 ③ Q-update ─────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "모델링 ③  ·  Q-update", "몬테카를로(MC) 방식 — 에피소드 종료 후 일괄 갱신")
pic_fit(s, fig_qflow(), Inches(0.6), Inches(1.75), Inches(11), Inches(2.0))
rect(s, Inches(2.8), Inches(3.95), Inches(6.6), Inches(0.85), LBLUE)
tf = add_textbox(s, Inches(2.9), Inches(4.02), Inches(6.4), Inches(0.7), MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "Q(s,a) ← Q(s,a) + α [ g − Q(s,a) ]      (α = 0.1)"
_set(r, 17, NAVY, bold=True); tf.paragraphs[0].alignment = PP_ALIGN.CENTER
bullets(s, [
    "왜 TD가 아닌 MC인가?",
    (1, "TD의 maximization bias(과대평가 편향) + 포커의 큰 분산 → 학습 붕괴"),
    (1, "MC는 실제 핸드 결과만 평균내므로 그 편향이 없음"),
], Inches(0.7), Inches(5.0), Inches(11), Inches(1.3), size=15, gap=6)

# ── 슬라이드 7: 모델링 - 탐색 ─────────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "모델링 ④  ·  탐색 (Exploration)", "온도 기반 Softmax — 탐색에서 활용으로")
pic_fit(s, fig_temp(), Inches(0.6), Inches(1.8), Inches(11), Inches(3.2))
bullets(s, [
    "P(a) = exp(Q(a)/τ) / Σ exp(Q(a')/τ)  —  Q값에 비례한 확률적 탐색",
    "온도 τ : 초반 10.0(탐색) → 후반 0.5(활용), 전체 80% 구간 선형 감소",
], Inches(0.7), Inches(5.15), Inches(11), Inches(1.2), size=15, gap=6)

# ── 슬라이드 8: 모델링 - 학습 데이터(핵심) ────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "모델링 ⑤  ·  학습 데이터 (제안의 핵심)", "알고리즘 고정, 매 에피소드 상대만 교체")
pic_fit(s, fig_personas(), Inches(0.55), Inches(1.75), Inches(7.0), Inches(4.4))
bullets(s, [
    "controlled comparison : 다른 조건 전부 고정 → '상대 다양화' 효과만 분리 측정",
    "두 가지 혼합 방식",
    (1, "cycle : 5종을 순서대로 균등 순환 (재현성)"),
    (1, "mixed : 가중 확률 샘플링, 약한 고리(LAG·STA) 비중↑"),
    "포지션(BB/SB)도 매 에피소드 번갈아 학습",
], Inches(7.7), Inches(1.95), Inches(4.1), Inches(4.2), size=14, gap=8)

# ── 슬라이드 9: 성능 평가 ─────────────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "성능 평가", "미학습 상대·정석 상대 양쪽에서 동시 흑자")
pic_fit(s, fig_results(), Inches(0.55), Inches(1.75), Inches(7.2), Inches(4.4))
bullets(s, [
    "평가 = 학습과 분리 : 탐색 끈 greedy 정책으로 대결",
    (1, "① 순수 Random  ② 학습에 안 쓴 고정 TAG"),
    "지표 mbb/g : 승률이 아닌 수익률",
    (1, "승률 높아도 큰 팟 잃으면 적자일 수 있음"),
    "규모 : 10만 게임 × 5회 반복",
    "합격선(사전등록) : 평균 − 1×표준편차 > 0",
    (1, "두 상대 모두 충족 ✅"),
], Inches(7.9), Inches(1.95), Inches(4.0), Inches(4.3), size=14, gap=7)

# ── 슬라이드 10: 마무리 ───────────────────────────────
s = prs.slides.add_slide(BLANK)
content_header(s, "마무리", "기여와 독창성")
bullets(s, [
    "해석 가능한 에이전트 : 신경망 없이 2,048칸 Q-테이블만으로 학습",
    (1, "어느 상황에서 무엇을 배웠는지 직접 들여다볼 수 있음"),
    "핵심 기여 : '상대 분포 다양화'가 과적합을 깨고 분포 강건성을 만든다",
    (1, "controlled comparison으로 두 독립 방식에서 입증"),
    "포커 특성에 맞춘 설계 3종",
    (1, "① 비율 기반 신용할당  ② MC로 TD 편향 회피  ③ 수익률(mbb/g) 평가"),
], Inches(0.7), Inches(1.9), Inches(11), Inches(4.0), size=16, gap=10)

# ── 슬라이드 11: 감사합니다 ───────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(3.95), SW, Pt(4), ACCENT)
tf = add_textbox(s, Inches(1.0), Inches(2.7), Inches(10.2), Inches(1.6), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run()
r.text = "감사합니다"; _set(r, 40, WHITE, bold=True)
tf2 = add_textbox(s, Inches(1.0), Inches(4.3), Inches(10.2), Inches(0.8), MSO_ANCHOR.MIDDLE)
p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run()
r.text = "Q & A"; _set(r, 20, RGBColor(0xBD, 0xD3, 0xEE))

OUT = os.path.join(os.path.dirname(__file__), "포커RL_발표.pptx")
try:
    prs.save(OUT)
except PermissionError:
    import time as _t
    OUT = os.path.join(os.path.dirname(__file__), f"포커RL_발표_{_t.strftime('%H%M%S')}.pptx")
    prs.save(OUT)
print("SAVED:", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides._sldIdLst), "slides")
